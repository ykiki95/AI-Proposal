"""
PPT 제안서 처리 모듈.
- python-pptx로 슬라이드 정보 추출
- 가로/세로 자동 판별 (슬라이드 너비/높이 비율)
- 슬라이드별 텍스트, 도형, 표 추출
- 재사용 가능한 슬라이드 패턴 자동 식별
"""

from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Any

from loguru import logger

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

# 슬라이드 패턴 키워드 매핑
_SLIDE_PATTERNS: dict[str, list[str]] = {
    "cover":            ["제안서", "제안", "사업명", "제목", "발표자료", "기술제안"],
    "toc":              ["목차", "TABLE OF CONTENTS", "INDEX", "차례", "CONTENTS"],
    "chapter_divider":  [],  # 텍스트 매우 짧은 슬라이드
    "as_is_to_be":      ["AS-IS", "TO-BE", "현황", "목표", "개선전", "개선후", "before", "after"],
    "architecture":     ["아키텍처", "architecture", "시스템 구성", "구성도", "인프라", "시스템도"],
    "schedule":         ["일정", "schedule", "wbs", "마일스톤", "milestone", "추진일정", "로드맵"],
    "org_chart":        ["조직", "인력", "팀 구성", "담당자", "투입인력", "수행조직", "pm"],
    "acceptance":       ["수용", "요구사항", "수용여부", "수용표", "요구조건"],
    "case_study":       ["사례", "실적", "레퍼런스", "reference", "수행실적", "경험"],
}


def _classify_slide_pattern(title: str, texts: list[str], slide_num: int) -> str:
    """슬라이드 내용으로 패턴 분류."""
    combined = (title + " " + " ".join(texts)).lower()
    word_count = len(combined.split())

    if slide_num <= 2 and word_count <= 30:
        return "cover"

    for pattern, keywords in _SLIDE_PATTERNS.items():
        if pattern == "chapter_divider":
            continue
        for kw in keywords:
            if kw.lower() in combined:
                return pattern

    if word_count <= 15 and slide_num > 2:
        return "chapter_divider"

    return "content"


class PPTXProcessor:
    """PPTX 제안서를 슬라이드 단위 JSON으로 변환하고 가로/세로 분류."""

    def __init__(self, processed_dir: Path | None = None, patterns_dir: Path | None = None):
        if not _PPTX_AVAILABLE:
            raise ImportError("python-pptx 미설치. pip install python-pptx")
        if processed_dir is None:
            from config.settings import PROCESSED_DIR
            processed_dir = PROCESSED_DIR
        if patterns_dir is None:
            from config.settings import PATTERNS_DIR
            patterns_dir = PATTERNS_DIR
        self.processed_dir = Path(processed_dir)
        self.patterns_dir = Path(patterns_dir)
        self.processed_dir.mkdir(parents=True, exist_ok=True)

    def process(self, pptx_path: str | Path) -> dict[str, Any]:
        """PPTX 파일 처리 → 구조화된 dict 반환 + JSON 저장."""
        pptx_path = Path(pptx_path)
        logger.info(f"PPTX 처리 시작: {pptx_path.name}")

        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logger.error(f"PPTX 열기 실패 ({pptx_path.name}): {e}")
            raise

        # 레이아웃 판별
        width = prs.slide_width
        height = prs.slide_height
        layout_mode = "horizontal" if width >= height else "vertical"

        slides_data: list[dict[str, Any]] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_info = self._extract_slide(slide, slide_num)
            slides_data.append(slide_info)

        # 패턴 통계
        pattern_counts: dict[str, int] = {}
        for s in slides_data:
            p = s["pattern"]
            pattern_counts[p] = pattern_counts.get(p, 0) + 1

        result: dict[str, Any] = {
            "filename": pptx_path.name,
            "layout_mode": layout_mode,
            "slide_dimensions": {
                "width_emu": width,
                "height_emu": height,
                "ratio": f"{width/height:.2f}",
            },
            "total_slides": len(slides_data),
            "pattern_summary": pattern_counts,
            "slides": slides_data,
        }

        # JSON 저장
        out_path = self.processed_dir / f"{pptx_path.stem}_pptx.json"
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)

        # 가로/세로 폴더에 메타데이터 기록
        self._record_in_pattern_folder(pptx_path, layout_mode, result)

        logger.info(
            f"PPTX 완료: {pptx_path.name} "
            f"[{layout_mode}] {len(slides_data)}슬라이드 → {out_path.name}"
        )
        return result

    def _extract_slide(self, slide: Any, slide_num: int) -> dict[str, Any]:
        """슬라이드에서 텍스트, 도형, 표 정보 추출."""
        title = ""
        texts: list[str] = []
        tables: list[list[list[str]]] = []
        shape_types: list[str] = []

        # 제목 추출
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            # 도형 유형 수집
            try:
                shape_types.append(shape.shape_type.name)
            except Exception:
                shape_types.append("UNKNOWN")

            # 텍스트 추출 (제목 제외)
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != title:
                        texts.append(text)

            # 표 추출
            if shape.has_table:
                table_data: list[list[str]] = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables.append(table_data)

        pattern = _classify_slide_pattern(title, texts, slide_num)

        return {
            "slide_number": slide_num,
            "title": title,
            "texts": texts,
            "tables": tables,
            "shape_types": list(set(shape_types)),
            "pattern": pattern,
        }

    def _record_in_pattern_folder(
        self, pptx_path: Path, layout_mode: str, result: dict[str, Any]
    ) -> None:
        """가로/세로 패턴 폴더에 메타데이터 JSON 기록."""
        target_dir = self.patterns_dir / layout_mode
        target_dir.mkdir(parents=True, exist_ok=True)
        meta_path = target_dir / f"{pptx_path.stem}_meta.json"
        meta = {
            "original_file": str(pptx_path),
            "filename": pptx_path.name,
            "layout_mode": layout_mode,
            "total_slides": result["total_slides"],
            "slide_dimensions": result["slide_dimensions"],
            "pattern_summary": result["pattern_summary"],
        }
        with open(meta_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)


def process_pptx(pptx_path: str | Path) -> dict[str, Any] | None:
    """단일 PPTX 처리 진입점."""
    try:
        return PPTXProcessor().process(pptx_path)
    except Exception as e:
        logger.error(f"PPTX 처리 실패 ({Path(pptx_path).name}): {e}")
        return None
