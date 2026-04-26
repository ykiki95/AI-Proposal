"""
박제안 기획 PPT(스토리보드) 생성기.

박제안이 작성한 본문 섹션을 받아 LLM이 슬라이드 단위 '기획안'으로 변환한다.
각 스토리보드 슬라이드는 다음 4개 박스로 시각화된다:
  ┌─────────── 핵심 메시지 ───────────┬─────────── 시각요소 ───────────┐
  │ (한 줄 헤드라인)                  │ (이미지/예시화면/표/그래프/도형)  │
  ├─────────── 레이아웃 ────────────┼─────────── 작성 프롬프트 ────────┤
  │ (페이지 안에서 무엇을 어디에)       │ (최종 PT 단계의 작성 지시문)     │
  └────────────────────────────────┴───────────────────────────────┘

슬라이드 노트(speaker notes)에는 RFP 매핑·전체 작성 지시문을 풀 텍스트로 적어
사장님이 슬라이드 한 장씩 검토하면서 "이 페이지가 어떻게 만들어질지"를 미리 확인할 수 있게 한다.
"""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import List, Optional

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import MODELS, OUTPUT_DIR, get_effective_company
from schemas.models import BidNotice, ProposalDraft
from tools.llm_clients import chat_anthropic
from tools.pptx_generator import _palette  # 디자인 팔레트 재사용

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x4B, 0x55, 0x63)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


SYSTEM_PROMPT = """당신은 한국 공공/민간 제안 PT를 100건 이상 기획한 시니어 발표 디렉터입니다.
박제안 팀이 작성한 본문 섹션을 받아, 최종 PT가 만들어지기 전에 사장이 검토할 수 있는
'슬라이드 단위 기획안(스토리보드)'으로 변환합니다.

각 본문 섹션은 평균 3~5장의 스토리보드 슬라이드로 분할합니다.
각 슬라이드는 RFP 목차를 90% 이상 준수해야 하며, 최종 PT에서 어떤 모습일지를 사장이
미리 머릿속에 그릴 수 있을 만큼 구체적으로 기획되어야 합니다.

출력은 JSON 배열만, 설명 금지."""


def _build_prompt(bid: BidNotice, draft: ProposalDraft, rfp_outline: str) -> str:
    company = get_effective_company()
    sections_text = "\n\n".join(
        f"## {sec.title}\n{sec.body[:2000]}"
        for sec in sorted(draft.sections, key=lambda s: s.order)
    )
    return f"""다음 제안서 본문을 슬라이드 단위 '기획안(스토리보드)'으로 변환하세요.

# 사업명
{bid.title}

# 제안 회사
{company.name} (대표 {company.ceo})

# RFP 목차 (가능하면 이 순서·번호를 슬라이드 rfp_section에 그대로 매핑하세요. 90% 이상 준수)
{rfp_outline}

# 박제안 팀 본문
{sections_text}

# 출력 형식 (JSON 배열만)
- 본문 섹션 8개 → 총 25~35장 권장 (각 섹션 평균 3~5장)
- 표지/목차/마무리 슬라이드는 제외 (별도로 자동 추가됨)
- 각 슬라이드의 visual_type은 다음 중 하나:
    "diagram"(도형/플로우), "screen"(예시 화면 mockup),
    "table"(표), "chart"(그래프), "image"(이미지),
    "metric"(정량 강조), "text"(텍스트 위주)

[
  {{
    "page": 1,
    "title": "사업 추진 배경",
    "rfp_section": "Ⅰ. 1. 사업 추진 배경",
    "key_message": "공공 디지털 서비스 수요 증가에 따라 시민 참여형 플랫폼 신규 구축이 필요하다.",
    "layout": "상단 1/3: 배경 통계 3개 카드 / 중단 1/3: As-Is 문제 도식 / 하단 1/3: 정책 인용 박스",
    "visual_type": "diagram",
    "visual_detail": "도형: 3단 카드 + As-Is/To-Be 화살표 / 색상: 네이비 + 오렌지 액센트",
    "specialist_prompt": "사업이해 전문가가 본문에서 다음을 슬라이드로 압축: 정책 배경 1줄, 통계 3개, As-Is 페인포인트 3개, 참고 인용 1개. 거짓 수치 금지."
  }}
]

스토리보드는 박제안의 본문이 갖춘 모든 핵심을 다 담아야 합니다 (시스템 아키텍처, 서버 아키텍처, WBS, 인력 R&R, 보안 통제, 가격 산정, 유사 실적 등).
"""


def _ask_llm(bid: BidNotice, draft: ProposalDraft) -> List[dict]:
    rfp_outline = (bid.rfp_full_text or bid.rfp_summary or "")[:3000]
    raw = chat_anthropic(
        MODELS.park,
        SYSTEM_PROMPT,
        _build_prompt(bid, draft, rfp_outline),
        max_tokens=8000,
    )
    m = re.search(r"\[.*\]", raw, re.S)
    if not m:
        raise ValueError("스토리보드 JSON 추출 실패")
    arr = json.loads(m.group(0))
    out = []
    for item in arr:
        if not isinstance(item, dict) or not item.get("title"):
            continue
        item.setdefault("rfp_section", "(매핑 없음)")
        item.setdefault("visual_type", "text")
        item.setdefault("key_message", "")
        item.setdefault("layout", "")
        item.setdefault("visual_detail", "")
        item.setdefault("specialist_prompt", "")
        out.append(item)
    return out


# ---------------------------------------------------------------------------
# 시각화
# ---------------------------------------------------------------------------

VISUAL_BADGE = {
    "diagram": ("◇ 도형/플로우", "diagram"),
    "screen": ("🖥 예시 화면", "screen"),
    "table": ("☷ 표", "table"),
    "chart": ("📊 그래프", "chart"),
    "image": ("🖼 이미지", "image"),
    "metric": ("📈 정량 강조", "metric"),
    "text": ("✎ 텍스트", "text"),
}


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _box(slide, x, y, w, h, fill, border=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(shp, fill)
    if border:
        shp.line.fill.solid()
        shp.line.fill.fore_color.rgb = border
    return shp


def _text(slide, left, top, width, height, text, *, size=11, bold=False,
          color=GRAY_DARK, align=PP_ALIGN.LEFT, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def _draw_storyboard_cover(prs, bid: BidNotice, total_slides: int, pal: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, pal["primary"])
    _text(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
          "📋 PROPOSAL STORYBOARD", size=18, bold=True, color=pal["warm"])
    _text(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6),
          bid.title, size=32, bold=True, color=WHITE)
    _text(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.5),
          f"발주기관 · {bid.agency}", size=14, color=pal["soft"])
    _text(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5),
          f"예상 슬라이드 · {total_slides}장 / 디자인 검토용 기획안", size=14, color=pal["soft"])
    _text(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
          "※ 이 문서는 최종 PT가 아닌 '슬라이드 기획안'입니다. "
          "각 슬라이드의 메시지·레이아웃·시각요소·작성 지시를 검토하시고 게이트2에서 승인해 주세요.",
          size=10, color=pal["soft"])


def _draw_storyboard_slide(prs, item: dict, idx: int, total: int, pal: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 배경
    _box(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    # 상단바 (제목 + RFP 매핑 + 페이지)
    _box(slide, 0, 0, SLIDE_W, Inches(0.85), pal["primary"])
    _text(slide, Inches(0.4), Inches(0.13), Inches(0.9), Inches(0.35),
          f"S{idx:02d}", size=14, bold=True, color=pal["warm"])
    _text(slide, Inches(1.3), Inches(0.13), Inches(8.5), Inches(0.35),
          item.get("title", ""), size=15, bold=True, color=WHITE)
    _text(slide, Inches(1.3), Inches(0.48), Inches(10.5), Inches(0.3),
          f"📍 RFP 매핑: {item.get('rfp_section', '')}",
          size=10, color=pal["soft"])
    _text(slide, SLIDE_W - Inches(1.5), Inches(0.25), Inches(1.3), Inches(0.4),
          f"{idx} / {total}", size=11, color=WHITE, align=PP_ALIGN.RIGHT)

    # 4분할 카드 영역
    inner_top = Inches(1.0)
    inner_h = SLIDE_H - inner_top - Inches(0.4)
    half_w = (SLIDE_W - Inches(0.6)) / 2
    half_h = (inner_h - Inches(0.15)) / 2
    pad = Inches(0.15)

    badge_label, _ = VISUAL_BADGE.get(item.get("visual_type", "text"), VISUAL_BADGE["text"])

    cards = [
        ("💬 핵심 메시지", item.get("key_message", ""), Inches(0.2), inner_top,
         half_w, half_h, pal["soft"], pal["primary"]),
        (f"{badge_label} · 시각요소", item.get("visual_detail", ""),
         Inches(0.2) + half_w + pad, inner_top, half_w - pad, half_h,
         pal["soft"], pal["accent"]),
        ("🧭 레이아웃 안내", item.get("layout", ""), Inches(0.2), inner_top + half_h + pad,
         half_w, half_h, WHITE, pal["primary"]),
        ("📝 작성 프롬프트", item.get("specialist_prompt", ""),
         Inches(0.2) + half_w + pad, inner_top + half_h + pad,
         half_w - pad, half_h, WHITE, pal["warm"]),
    ]
    for label, body, x, y, w, h, fill, header_color in cards:
        # 카드 배경
        _box(slide, x, y, w, h, fill, border=header_color)
        # 카드 헤더
        _box(slide, x, y, w, Inches(0.38), header_color)
        _text(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.3),
              label, size=11, bold=True, color=WHITE)
        # 카드 본문
        _text(slide, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3), h - Inches(0.55),
              body or "(미작성)", size=11, color=GRAY_DARK)

    # 슬라이드 노트에 풀 텍스트
    notes = (
        f"[페이지 {idx}] {item.get('title','')}\n"
        f"RFP 매핑: {item.get('rfp_section','')}\n"
        f"시각요소 종류: {item.get('visual_type','')}\n\n"
        f"■ 핵심 메시지\n{item.get('key_message','')}\n\n"
        f"■ 레이아웃 안내\n{item.get('layout','')}\n\n"
        f"■ 시각요소 상세\n{item.get('visual_detail','')}\n\n"
        f"■ 최종 PT 작성 프롬프트\n{item.get('specialist_prompt','')}\n"
    )
    slide.notes_slide.notes_text_frame.text = notes


def _draw_storyboard_summary(prs, items: List[dict], pal: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _box(slide, 0, 0, SLIDE_W, Inches(0.85), pal["primary"])
    _text(slide, Inches(0.5), Inches(0.22), Inches(12.0), Inches(0.55),
          "📑 STORYBOARD INDEX", size=20, bold=True, color=WHITE)
    # 페이지·제목·시각요소 목록
    box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items, start=1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        badge, _ = VISUAL_BADGE.get(it.get("visual_type", "text"), VISUAL_BADGE["text"])
        r1 = p.add_run()
        r1.text = f"S{i:02d}  "
        r1.font.size = Pt(10)
        r1.font.bold = True
        r1.font.color.rgb = pal["warm"]
        r2 = p.add_run()
        r2.text = f"{it.get('title','')}  "
        r2.font.size = Pt(11)
        r2.font.color.rgb = GRAY_DARK
        r3 = p.add_run()
        r3.text = f"  [{badge}]  "
        r3.font.size = Pt(9)
        r3.font.color.rgb = pal["accent"]
        r4 = p.add_run()
        r4.text = f"  📍{it.get('rfp_section','')}"
        r4.font.size = Pt(9)
        r4.font.color.rgb = GRAY
        p.space_after = Pt(2)


# ---------------------------------------------------------------------------
# 진입점
# ---------------------------------------------------------------------------

def render_storyboard(bid: BidNotice, draft: ProposalDraft,
                      template: Optional[str] = None) -> Path:
    """박제안 본문을 받아 기획 PPT를 생성하고 경로를 반환."""
    if template is None:
        try:
            from tools.db import get_design_template
            template = get_design_template()
        except Exception:
            template = "navy"
    pal = _palette(template)

    items = _ask_llm(bid, draft)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _draw_storyboard_cover(prs, bid, len(items), pal)
    _draw_storyboard_summary(prs, items, pal)
    for i, item in enumerate(items, start=1):
        _draw_storyboard_slide(prs, item, i, len(items), pal)

    fname = f"{bid.bid_id}_storyboard_v{draft.version}.pptx"
    out_path = OUTPUT_DIR / fname
    prs.save(out_path)
    logger.info(f"📋 스토리보드 PPTX 생성: {out_path} (슬라이드 {len(items)}장)")
    return out_path
