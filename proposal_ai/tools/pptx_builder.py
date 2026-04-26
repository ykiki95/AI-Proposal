"""PPTX 빌더 v2 — 박제안 팀 slides_json + 정디자 DesignBrief → 최종 .pptx.

핵심:
- DesignBrief.master_path 가 있으면 그 .pptx 를 베이스로 로드 (사용자 등록 마스터).
- 없으면 빈 Presentation에 Theme 색으로 그리기.
- DesignBrief.accent_hex 가 있으면 Theme.accent 를 오버라이드.
- 슬라이드 layout 종류:
    cover, toc, section_divider, closing,
    title_bullets, two_column_compare, diagram_layered,
    metric_cards, table, blank_placeholder,
    as_is_to_be_compare, system_architecture, process_flow,
    screen_mockup, screen_mockup_grid
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import List, Optional

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import OUTPUT_DIR, get_effective_company
from schemas.models import BidNotice, DesignBrief, ProposalDraft
from tools.pptx_themes import THEMES, Theme, get_theme, hex_to_rgb

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, x, y, w, h, text, *, size=14, bold=False, color=None,
          align=PP_ALIGN.LEFT, font="맑은 고딕") -> None:
    if color is None:
        color = RGBColor(0x1F, 0x29, 0x37)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def _theme_with_accent(brief: DesignBrief) -> Theme:
    """DesignBrief.accent_hex 가 있으면 Theme.accent 를 오버라이드한 새 Theme 반환."""
    base = get_theme(brief.theme_key)
    rgb = hex_to_rgb(brief.accent_hex)
    if not rgb:
        return base
    return Theme(
        key=base.key, label=base.label, tone=base.tone,
        primary=base.primary, accent=rgb, warm=base.warm, soft=base.soft,
        text_dark=base.text_dark, text_gray=base.text_gray,
        font_title=base.font_title, font_body=base.font_body,
    )


def _new_presentation(brief: DesignBrief) -> Presentation:
    """사용자 마스터 .pptx 가 있으면 그걸 기반으로, 없으면 빈 Presentation."""
    if brief.master_path and Path(brief.master_path).exists():
        try:
            prs = Presentation(brief.master_path)
            # 사용자 마스터는 미리보기 슬라이드를 포함할 수 있으므로 기존 슬라이드는 제거
            xml_slides = prs.slides._sldIdLst
            for sld in list(xml_slides):
                xml_slides.remove(sld)
            prs.slide_width = SLIDE_W
            prs.slide_height = SLIDE_H
            logger.info(f"마스터 베이스 사용: {brief.master_path}")
            return prs
        except Exception as e:
            logger.warning(f"마스터 로드 실패 ({brief.master_path}) → 빈 Presentation: {e}")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --- 슬라이드 렌더러 -------------------------------------------------

def _draw_cover(prs, theme: Theme, bid: BidNotice, company) -> None:
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, theme.rgb("primary"))
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.7), 0, Inches(1.63), SLIDE_H)
    _solid(accent, theme.rgb("accent"))
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.3), Inches(0.12), Inches(2.3))
    _solid(line, theme.rgb("warm"))
    _text(s, Inches(0.85), Inches(2.0), Inches(10.0), Inches(0.5),
          "제 안 서", size=16, bold=True, color=theme.rgb("soft"), font=theme.font_title)
    _text(s, Inches(0.85), Inches(2.55), Inches(11.0), Inches(2.0),
          bid.title, size=32, bold=True, color=WHITE, font=theme.font_title)
    _text(s, Inches(0.85), Inches(5.3), Inches(10.0), Inches(0.4),
          f"발주기관  |  {bid.agency}", size=13, color=theme.rgb("soft"), font=theme.font_body)
    _text(s, Inches(0.85), Inches(5.8), Inches(10.0), Inches(0.4),
          f"제안사  |  {company.name}  (대표 {company.ceo})",
          size=13, color=WHITE, font=theme.font_body)
    _text(s, Inches(0.85), Inches(6.7), Inches(10.0), Inches(0.4),
          f"제출일자  |  {datetime.now().strftime('%Y년 %m월 %d일')}",
          size=11, color=theme.rgb("soft"), font=theme.font_body)


def _draw_toc(prs, theme: Theme, sections) -> None:
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, WHITE)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    _solid(band, theme.rgb("primary"))
    _text(s, Inches(0.5), Inches(0.22), Inches(12.0), Inches(0.55),
          "C O N T E N T S", size=22, bold=True, color=WHITE, font=theme.font_title)
    titles = [sec.title for sec in sorted(sections, key=lambda x: x.order)]
    half = (len(titles) + 1) // 2
    for col, sub in enumerate([titles[:half], titles[half:]]):
        for i, t in enumerate(sub):
            idx = col * half + i + 1
            x = Inches(0.7 + col * 6.3)
            y = Inches(1.3 + i * 0.55)
            circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
            _solid(circ, theme.rgb("accent"))
            tf = circ.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{idx:02d}"
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = WHITE
            _text(s, x + Inches(0.55), y + Inches(0.04), Inches(5.7), Inches(0.4),
                  t, size=14, bold=True, color=theme.rgb("text_dark"), font=theme.font_title)


def _title_bar(s, theme: Theme, idx: int, title: str) -> None:
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SLIDE_H)
    _solid(bar, theme.rgb("accent"))
    title_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.22), 0, SLIDE_W - Inches(0.22), Inches(0.75))
    _solid(title_bg, theme.rgb("soft"))
    _text(s, Inches(0.45), Inches(0.18), Inches(0.7), Inches(0.45),
          f"{idx:02d}", size=16, bold=True, color=theme.rgb("accent"), font=theme.font_title)
    _text(s, Inches(1.1), Inches(0.2), Inches(11.5), Inches(0.5),
          title, size=18, bold=True, color=theme.rgb("primary"), font=theme.font_title)


def _footer(s, theme: Theme, footer: str, idx: int, total: int) -> None:
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3))
    _solid(bar, theme.rgb("primary"))
    _text(s, Inches(0.4), SLIDE_H - Inches(0.28), Inches(10.0), Inches(0.26),
          footer, size=9, color=WHITE, font=theme.font_body)
    _text(s, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.28), Inches(1.0), Inches(0.26),
          f"{idx} / {total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT, font=theme.font_body)


def _draw_section_divider(s, theme: Theme, idx: int, title: str, subtitle: str) -> None:
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, theme.rgb("primary"))
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), Inches(1.2), Inches(0.12))
    _solid(line, theme.rgb("warm"))
    _text(s, Inches(1.4), Inches(2.7), Inches(11.0), Inches(0.5),
          f"PART  {idx:02d}", size=14, bold=True,
          color=theme.rgb("warm"), font=theme.font_title)
    _text(s, Inches(1.4), Inches(3.3), Inches(11.0), Inches(1.4),
          title, size=40, bold=True, color=WHITE, font=theme.font_title)
    if subtitle:
        _text(s, Inches(1.4), Inches(4.7), Inches(11.0), Inches(0.7),
              subtitle, size=16, color=theme.rgb("soft"), font=theme.font_body)


def _draw_title_bullets(s, theme: Theme, bullets: List[str]) -> None:
    box = s.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.4), Inches(5.9))
    tf = box.text_frame
    tf.word_wrap = True
    items = bullets or ["(불릿 없음)"]
    for i, line in enumerate(items[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        m = p.add_run()
        m.text = "▪  "
        m.font.size = Pt(14)
        m.font.color.rgb = theme.rgb("warm")
        m.font.bold = True
        r = p.add_run()
        r.text = str(line)
        r.font.name = theme.font_body
        r.font.size = Pt(14)
        r.font.color.rgb = theme.rgb("text_dark")
        p.space_after = Pt(8)


def _draw_two_column(s, theme: Theme, left: dict, right: dict) -> None:
    col_w = Inches(6.0)
    col_h = Inches(5.6)
    top = Inches(1.1)
    for i, (col, color) in enumerate([
        (left or {}, theme.rgb("text_gray")),
        (right or {}, theme.rgb("accent")),
    ]):
        x = Inches(0.55) + (col_w + Inches(0.3)) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _solid(card, theme.rgb("soft") if i == 0 else WHITE)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.55))
        _solid(head, color)
        _text(s, x + Inches(0.25), top + Inches(0.1), col_w, Inches(0.4),
              col.get("label", "AS-IS" if i == 0 else "TO-BE"),
              size=14, bold=True, color=WHITE, font=theme.font_title)
        items = col.get("items", []) or ["(내용 없음)"]
        box = s.shapes.add_textbox(x + Inches(0.3), top + Inches(0.8),
                                   col_w - Inches(0.6), col_h - Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        for j, it in enumerate(items[:8]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            m = p.add_run()
            m.text = "•  "
            m.font.color.rgb = color
            m.font.bold = True
            r = p.add_run()
            r.text = str(it)
            r.font.name = theme.font_body
            r.font.size = Pt(13)
            r.font.color.rgb = theme.rgb("text_dark")
            p.space_after = Pt(6)


def _draw_diagram_layered(s, theme: Theme, diagram: dict) -> None:
    layers = (diagram or {}).get("layers", [])
    if not layers:
        layers = [{"name": "(레이어 없음)", "items": []}]
    layers = layers[:5]
    top = Inches(1.1)
    bottom = SLIDE_H - Inches(0.5)
    avail = bottom - top
    n = len(layers)
    layer_h = avail / n
    gap = Inches(0.1)
    for li, layer in enumerate(layers):
        y = top + layer_h * li + gap
        h = layer_h - gap
        # 레이어 라벨 박스
        lbl_w = Inches(2.0)
        lbl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, lbl_w, h)
        _solid(lbl, theme.rgb("primary"))
        _text(s, Inches(0.55), y + Inches(0.15), lbl_w, h - Inches(0.3),
              str(layer.get("name", "")), size=13, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, font=theme.font_title)
        # 아이템 박스들
        items = (layer.get("items", []) or [""])[:6]
        item_area_w = SLIDE_W - Inches(0.55) - lbl_w - Inches(0.55)
        item_w = item_area_w / len(items)
        for ii, it in enumerate(items):
            x = Inches(0.55) + lbl_w + Inches(0.1) + (item_w - Inches(0.1)) * 0 + (item_area_w / len(items)) * ii + Inches(0.05)
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x, y, item_w - Inches(0.1), h)
            _solid(box, theme.rgb("soft"))
            _text(s, x, y + Inches(0.15), item_w - Inches(0.1), h - Inches(0.3),
                  str(it), size=11, color=theme.rgb("text_dark"),
                  align=PP_ALIGN.CENTER, font=theme.font_body)


def _draw_metric_cards(s, theme: Theme, metrics: list) -> None:
    if not metrics:
        metrics = [{"value": "—", "label": "(수치 없음)"}]
    metrics = metrics[:4]
    n = len(metrics)
    gap = Inches(0.25)
    card_w = Inches((12.4 - 0.25 * (n - 1)) / n)
    top = Inches(1.6)
    h = Inches(4.2)
    for i, m in enumerate(metrics):
        x = Inches(0.55) + (card_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, h)
        _solid(card, theme.rgb("soft"))
        _text(s, x, top + Inches(0.7), card_w, Inches(1.5),
              m.get("value", ""), size=44, bold=True,
              color=theme.rgb("primary"), align=PP_ALIGN.CENTER, font=theme.font_title)
        _text(s, x, top + Inches(2.4), card_w, Inches(0.5),
              m.get("label", ""), size=14, color=theme.rgb("text_dark"),
              align=PP_ALIGN.CENTER, font=theme.font_body)
        if m.get("note"):
            _text(s, x, top + Inches(3.1), card_w, Inches(0.5),
                  m["note"], size=10, color=theme.rgb("text_gray"),
                  align=PP_ALIGN.CENTER, font=theme.font_body)


def _draw_table(s, theme: Theme, table_data: dict) -> None:
    headers = (table_data or {}).get("headers", [])
    rows = (table_data or {}).get("rows", [])
    if not headers and not rows:
        _text(s, Inches(0.55), Inches(2.5), Inches(12.4), Inches(0.5),
              "(표 데이터 없음)", size=14, color=theme.rgb("text_gray"),
              align=PP_ALIGN.CENTER, font=theme.font_body)
        return
    if not headers:
        headers = [f"열 {i+1}" for i in range(len(rows[0]) if rows else 1)]
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    left, top = Inches(0.55), Inches(1.05)
    width, height = Inches(12.4), Inches(5.8)
    shp = s.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for ci, h in enumerate(headers):
        cell = shp.cell(0, ci)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.rgb("primary")
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = theme.font_title
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for ri, row in enumerate(rows, start=1):
        for ci in range(n_cols):
            cell = shp.cell(ri, ci)
            cell.text = str(row[ci]) if ci < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.rgb("soft") if ri % 2 == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = theme.font_body
                    r.font.size = Pt(11)
                    r.font.color.rgb = theme.rgb("text_dark")


def _draw_as_is_to_be_compare(s, theme: Theme, as_is: List[str], to_be: List[str],
                              arrow_label: Optional[str]) -> None:
    """좌(현행 박스 5개) → 화살표(중앙) → 우(개선 박스 5개) 도형."""
    top = Inches(1.15)
    bottom = SLIDE_H - Inches(0.5)
    avail_h = bottom - top
    side_w = Inches(5.4)
    arrow_w = Inches(2.0)
    left_x = Inches(0.4)
    right_x = SLIDE_W - side_w - Inches(0.4)
    arrow_x = left_x + side_w + Inches(0.1)

    # 헤더 라벨 (AS-IS / TO-BE)
    head_h = Inches(0.55)
    head_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top, side_w, head_h)
    _solid(head_l, theme.rgb("text_gray"))
    _text(s, left_x, top + Inches(0.08), side_w, head_h,
          "AS-IS · 현행", size=15, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, font=theme.font_title)
    head_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top, side_w, head_h)
    _solid(head_r, theme.rgb("accent"))
    _text(s, right_x, top + Inches(0.08), side_w, head_h,
          "TO-BE · 개선 후", size=15, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, font=theme.font_title)

    # 박스들
    items_top = top + head_h + Inches(0.1)
    items_h = avail_h - head_h - Inches(0.1)
    as_is_list = (as_is or ["(현행 정보 없음)"])[:5]
    to_be_list = (to_be or ["(개선안 없음)"])[:5]
    n = max(len(as_is_list), len(to_be_list))
    box_h = (items_h - Inches(0.1) * (n - 1)) / max(n, 1)
    for i in range(n):
        y = items_top + (box_h + Inches(0.1)) * i
        # 좌측 (붉은 톤)
        if i < len(as_is_list):
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y, side_w, box_h)
            _solid(box, RGBColor(0xFD, 0xEC, 0xEA))
            box.line.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
            _text(s, left_x + Inches(0.2), y + Inches(0.1), side_w - Inches(0.4), box_h - Inches(0.2),
                  f"❌  {as_is_list[i]}", size=12, color=RGBColor(0x80, 0x21, 0x18),
                  font=theme.font_body)
        # 우측 (강조 톤)
        if i < len(to_be_list):
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, side_w, box_h)
            _solid(box, theme.rgb("soft"))
            box.line.color.rgb = theme.rgb("accent")
            _text(s, right_x + Inches(0.2), y + Inches(0.1), side_w - Inches(0.4), box_h - Inches(0.2),
                  f"✅  {to_be_list[i]}", size=12, bold=True,
                  color=theme.rgb("primary"), font=theme.font_body)

    # 중앙 화살표
    arr_y = top + avail_h / 2 - Inches(0.4)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arr_y, arrow_w, Inches(0.8))
    _solid(arrow, theme.rgb("warm"))
    if arrow_label:
        _text(s, arrow_x, arr_y - Inches(0.45), arrow_w, Inches(0.4),
              str(arrow_label), size=11, bold=True,
              color=theme.rgb("primary"), align=PP_ALIGN.CENTER, font=theme.font_body)


def _draw_system_architecture(s, theme: Theme, arch: dict) -> None:
    """레이어형 시스템 구성도 + 외부 연계.
    arch = {"layers":[{"name":"...", "items":["...","..."]}], "external":["기관A","기관B"]}
    """
    layers = (arch or {}).get("layers", []) or [{"name": "(레이어 없음)", "items": []}]
    layers = layers[:6]
    external = (arch or {}).get("external", []) or []
    external = external[:6]

    top = Inches(1.1)
    bottom = SLIDE_H - Inches(0.5)
    avail_h = bottom - top
    main_w = SLIDE_W - Inches(0.55) - Inches(0.55)
    if external:
        main_w -= Inches(2.6)
    n = len(layers)
    layer_h = (avail_h - Inches(0.1) * (n - 1)) / n

    # 메인 레이어 박스들
    for li, layer in enumerate(layers):
        y = top + (layer_h + Inches(0.1)) * li
        lbl_w = Inches(1.7)
        lbl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, lbl_w, layer_h)
        _solid(lbl, theme.rgb("primary"))
        _text(s, Inches(0.55), y + (layer_h - Inches(0.3)) / 2,
              lbl_w, Inches(0.5), str(layer.get("name", "")),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=theme.font_title)
        items = (layer.get("items", []) or [""])[:6]
        items_x = Inches(0.55) + lbl_w + Inches(0.1)
        items_w = main_w - lbl_w - Inches(0.1)
        per_w = (items_w - Inches(0.1) * (len(items) - 1)) / max(len(items), 1)
        for ii, it in enumerate(items):
            x = items_x + (per_w + Inches(0.1)) * ii
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, per_w, layer_h)
            _solid(box, theme.rgb("soft"))
            box.line.color.rgb = theme.rgb("primary")
            _text(s, x + Inches(0.1), y + Inches(0.15), per_w - Inches(0.2), layer_h - Inches(0.3),
                  str(it), size=11, color=theme.rgb("text_dark"),
                  align=PP_ALIGN.CENTER, font=theme.font_body)

    # 외부 연계 패널 (오른쪽)
    if external:
        ext_x = Inches(0.55) + main_w + Inches(0.15)
        ext_w = SLIDE_W - ext_x - Inches(0.4)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ext_x, top, ext_w, Inches(0.5))
        _solid(head, theme.rgb("warm"))
        _text(s, ext_x, top + Inches(0.06), ext_w, Inches(0.4),
              "외부 연계", size=12, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, font=theme.font_title)
        per_h = (avail_h - Inches(0.5) - Inches(0.1) * (len(external) - 1)) / max(len(external), 1)
        for ei, ex in enumerate(external):
            ey = top + Inches(0.6) + (per_h + Inches(0.1)) * ei
            box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ext_x, ey, ext_w, per_h)
            _solid(box, WHITE)
            box.line.color.rgb = theme.rgb("warm")
            _text(s, ext_x + Inches(0.05), ey + (per_h - Inches(0.3)) / 2,
                  ext_w - Inches(0.1), Inches(0.5), str(ex),
                  size=11, color=theme.rgb("text_dark"),
                  align=PP_ALIGN.CENTER, font=theme.font_body)


def _draw_process_flow(s, theme: Theme, steps: List[dict]) -> None:
    """가로 박스 + 화살표 chain."""
    if not steps:
        steps = [{"step": "1", "title": "(단계 없음)", "desc": ""}]
    steps = steps[:6]
    n = len(steps)
    top = Inches(2.5)
    box_h = Inches(2.4)
    pad = Inches(0.4)
    avail_w = SLIDE_W - pad * 2
    arrow_w = Inches(0.45) if n > 1 else Inches(0)
    box_w = (avail_w - arrow_w * (n - 1)) / n
    for i, st in enumerate(steps):
        x = pad + (box_w + arrow_w) * i
        # 단계 번호 (원)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + box_w / 2 - Inches(0.4), top - Inches(0.5),
                                  Inches(0.8), Inches(0.8))
        _solid(circ, theme.rgb("accent"))
        tf = circ.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(st.get("step") or str(i + 1))
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        # 본 박스
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        _solid(box, theme.rgb("soft"))
        box.line.color.rgb = theme.rgb("primary")
        _text(s, x + Inches(0.15), top + Inches(0.5), box_w - Inches(0.3), Inches(0.6),
              str(st.get("title", "")), size=14, bold=True,
              color=theme.rgb("primary"), align=PP_ALIGN.CENTER, font=theme.font_title)
        desc = str(st.get("desc", ""))
        if desc:
            _text(s, x + Inches(0.2), top + Inches(1.1), box_w - Inches(0.4), box_h - Inches(1.3),
                  desc, size=11, color=theme.rgb("text_dark"),
                  align=PP_ALIGN.CENTER, font=theme.font_body)
        # 화살표
        if i < n - 1:
            ay = top + box_h / 2 - Inches(0.2)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.05),
                                     ay, arrow_w - Inches(0.1), Inches(0.4))
            _solid(arr, theme.rgb("accent"))


def _draw_screen_mockup_one(s, theme: Theme, mockup: dict, x, y, w, h,
                             *, big: bool = True) -> None:
    """단일 화면 mockup (PC 또는 모바일 프레임). 영역(regions)별 박스로 표현."""
    device = (mockup or {}).get("device", "pc")  # pc | mobile
    header = (mockup or {}).get("header", "")
    footer = (mockup or {}).get("footer", "")
    regions = (mockup or {}).get("regions", []) or []
    label = (mockup or {}).get("label") or (mockup or {}).get("title") or ""

    # 라벨 (mockup 위)
    lbl_h = Inches(0.35) if big else Inches(0.28)
    if label:
        _text(s, x, y, w, lbl_h, label, size=12 if big else 10, bold=True,
              color=theme.rgb("primary"), align=PP_ALIGN.CENTER, font=theme.font_title)
    frame_y = y + (lbl_h if label else Inches(0))
    frame_h = h - (lbl_h if label else Inches(0))

    # 디바이스 프레임
    if device == "mobile":
        # 모바일: 16:9 세로 → 폭을 줄여서 중앙 배치
        target_w = frame_h * 9 / 16
        if target_w > w:
            target_w = w * 0.85
            frame_h = target_w * 16 / 9
        fx = x + (w - target_w) / 2
        fw = target_w
    else:
        fx = x
        fw = w

    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, frame_y, fw, frame_h)
    _solid(frame, RGBColor(0x21, 0x29, 0x37))
    frame.line.fill.background()

    # 내부 화면 영역 (검은 프레임 안의 흰 화면)
    bezel = Inches(0.08)
    sx = fx + bezel
    sy = frame_y + bezel
    sw = fw - bezel * 2
    sh = frame_h - bezel * 2
    screen = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, sw, sh)
    _solid(screen, WHITE)

    # 헤더 영역
    head_h = Inches(0.35) if big else Inches(0.22)
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy, sw, head_h)
    _solid(head, theme.rgb("primary"))
    _text(s, sx + Inches(0.1), sy + Inches(0.04), sw - Inches(0.2), head_h,
          str(header) or "헤더 / 로고", size=10 if big else 8,
          bold=True, color=WHITE, font=theme.font_title)

    # 본문 영역들 (regions를 세로로 분할)
    body_top = sy + head_h + Inches(0.05)
    body_h = sh - head_h - (Inches(0.3) if footer else Inches(0.1))
    if not regions:
        regions = [{"label": "콘텐츠 영역", "kind": "panel"}]
    regions = regions[:6]
    per_h = (body_h - Inches(0.05) * (len(regions) - 1)) / len(regions)
    for ri, rg in enumerate(regions):
        ry = body_top + (per_h + Inches(0.05)) * ri
        kind = (rg or {}).get("kind", "panel")  # panel | grid | list | hero | table | chart
        rlabel = (rg or {}).get("label", "")
        col = {
            "hero": theme.rgb("accent"),
            "table": theme.rgb("text_gray"),
            "chart": theme.rgb("warm"),
        }.get(kind, theme.rgb("soft"))
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx + Inches(0.08), ry,
                                 sw - Inches(0.16), per_h)
        _solid(box, col)
        # 텍스트 (다른 색)
        txt_color = WHITE if kind in ("hero", "table") else theme.rgb("text_dark")
        _text(s, sx + Inches(0.18), ry + Inches(0.05), sw - Inches(0.36), per_h,
              str(rlabel), size=10 if big else 8, bold=True,
              color=txt_color, font=theme.font_body)
        # grid이면 안에 작은 박스 4개
        if kind == "grid" and per_h > Inches(0.6):
            grid_top = ry + Inches(0.3)
            grid_h = per_h - Inches(0.4)
            cell_w = (sw - Inches(0.36) - Inches(0.05) * 3) / 4
            for ci in range(4):
                cx = sx + Inches(0.18) + (cell_w + Inches(0.05)) * ci
                cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, grid_top, cell_w, grid_h)
                _solid(cell, WHITE)
                cell.line.color.rgb = theme.rgb("text_gray")

    # 푸터
    if footer:
        ft = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, sy + sh - Inches(0.25),
                                 sw, Inches(0.25))
        _solid(ft, theme.rgb("text_gray"))
        _text(s, sx + Inches(0.1), sy + sh - Inches(0.23), sw - Inches(0.2), Inches(0.22),
              str(footer), size=8, color=WHITE, font=theme.font_body)


def _draw_screen_mockup(s, theme: Theme, mockup: dict) -> None:
    """단일 화면 mockup (큰 사이즈)."""
    pad_x = Inches(2.0)
    pad_y = Inches(1.05)
    w = SLIDE_W - pad_x * 2
    h = SLIDE_H - pad_y - Inches(0.5)
    _draw_screen_mockup_one(s, theme, mockup, pad_x, pad_y, w, h, big=True)


def _draw_screen_mockup_grid(s, theme: Theme, mockups: List[dict]) -> None:
    """4~6장 화면 mockup 그리드."""
    if not mockups:
        mockups = [{"label": "(화면 없음)", "regions": []}]
    mockups = mockups[:6]
    n = len(mockups)
    cols = 3 if n > 4 else 2
    rows = (n + cols - 1) // cols
    pad_x = Inches(0.45)
    pad_y = Inches(1.05)
    gap = Inches(0.2)
    avail_w = SLIDE_W - pad_x * 2 - gap * (cols - 1)
    avail_h = SLIDE_H - pad_y - Inches(0.5) - gap * (rows - 1)
    cell_w = avail_w / cols
    cell_h = avail_h / rows
    for i, mk in enumerate(mockups):
        r = i // cols
        c = i % cols
        x = pad_x + (cell_w + gap) * c
        y = pad_y + (cell_h + gap) * r
        _draw_screen_mockup_one(s, theme, mk, x, y, cell_w, cell_h, big=False)


def _draw_blank_placeholder(s, theme: Theme, title: str, reason: Optional[str]) -> None:
    pad = Inches(0.55)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             pad, Inches(1.1), SLIDE_W - pad * 2, Inches(5.8))
    _solid(box, theme.rgb("soft"))
    box.line.fill.solid()
    box.line.fill.fore_color.rgb = theme.rgb("accent")
    _text(s, pad, Inches(2.6), SLIDE_W - pad * 2, Inches(0.7),
          "👤  사람 작성 영역", size=22, bold=True,
          color=theme.rgb("primary"), align=PP_ALIGN.CENTER, font=theme.font_title)
    _text(s, pad, Inches(3.6), SLIDE_W - pad * 2, Inches(0.6),
          title, size=16, color=theme.rgb("text_dark"),
          align=PP_ALIGN.CENTER, font=theme.font_body)
    if reason:
        _text(s, pad, Inches(4.4), SLIDE_W - pad * 2, Inches(0.6),
              reason, size=12, color=theme.rgb("text_gray"),
              align=PP_ALIGN.CENTER, font=theme.font_body)


def _draw_closing(s, theme: Theme, company) -> None:
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, theme.rgb("primary"))
    _text(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.4),
          "감사합니다", size=52, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER, font=theme.font_title)
    _text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
          "함께 만들어 가겠습니다", size=18,
          color=theme.rgb("soft"), align=PP_ALIGN.CENTER, font=theme.font_body)
    _text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
          f"{company.name}  |  대표 {company.ceo}", size=12,
          color=theme.rgb("soft"), align=PP_ALIGN.CENTER, font=theme.font_body)


# --- 진입점 ---------------------------------------------------------

def _render_section_slides(prs, theme: Theme, section, footer: str,
                           start_idx: int, total: int) -> int:
    """섹션의 slides_json 모두 렌더링. 다음 슬라이드 인덱스 반환."""
    idx = start_idx
    for slide_dict in (section.slides_json or []):
        layout = slide_dict.get("layout", "title_bullets")
        title = slide_dict.get("title", section.title)
        s = _blank(prs)
        if layout == "section_divider":
            _draw_section_divider(s, theme, idx, title, slide_dict.get("subtitle", ""))
        elif layout == "blank_placeholder":
            _title_bar(s, theme, idx, title)
            _draw_blank_placeholder(s, theme, title, slide_dict.get("placeholder_reason"))
            _footer(s, theme, footer, idx, total)
        elif layout == "two_column_compare":
            _title_bar(s, theme, idx, title)
            _draw_two_column(s, theme, slide_dict.get("left"), slide_dict.get("right"))
            _footer(s, theme, footer, idx, total)
        elif layout == "diagram_layered":
            _title_bar(s, theme, idx, title)
            _draw_diagram_layered(s, theme, slide_dict.get("diagram") or {})
            _footer(s, theme, footer, idx, total)
        elif layout == "metric_cards":
            _title_bar(s, theme, idx, title)
            _draw_metric_cards(s, theme, slide_dict.get("metrics") or [])
            _footer(s, theme, footer, idx, total)
        elif layout == "table":
            _title_bar(s, theme, idx, title)
            _draw_table(s, theme, slide_dict.get("table_data") or {})
            _footer(s, theme, footer, idx, total)
        elif layout == "as_is_to_be_compare":
            _title_bar(s, theme, idx, title)
            _draw_as_is_to_be_compare(
                s, theme,
                slide_dict.get("as_is_items") or [],
                slide_dict.get("to_be_items") or [],
                slide_dict.get("arrow_label"),
            )
            _footer(s, theme, footer, idx, total)
        elif layout == "system_architecture":
            _title_bar(s, theme, idx, title)
            _draw_system_architecture(s, theme, slide_dict.get("architecture") or {})
            _footer(s, theme, footer, idx, total)
        elif layout == "process_flow":
            _title_bar(s, theme, idx, title)
            _draw_process_flow(s, theme, slide_dict.get("flow_steps") or [])
            _footer(s, theme, footer, idx, total)
        elif layout == "screen_mockup":
            _title_bar(s, theme, idx, title)
            _draw_screen_mockup(s, theme, slide_dict.get("mockup") or {})
            _footer(s, theme, footer, idx, total)
        elif layout == "screen_mockup_grid":
            _title_bar(s, theme, idx, title)
            _draw_screen_mockup_grid(s, theme, slide_dict.get("mockups") or [])
            _footer(s, theme, footer, idx, total)
        else:  # title_bullets / unknown
            _title_bar(s, theme, idx, title)
            _draw_title_bullets(s, theme, slide_dict.get("bullets") or [])
            _footer(s, theme, footer, idx, total)
        # 발표자 노트
        notes = slide_dict.get("speaker_notes") or ""
        if notes:
            try:
                s.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
        idx += 1
    return idx


def build_proposal_pptx_v2(bid: BidNotice, draft: ProposalDraft,
                            brief: DesignBrief) -> Path:
    theme = _theme_with_accent(brief)
    company = get_effective_company()
    prs = _new_presentation(brief)

    sorted_sections = sorted(draft.sections, key=lambda s: s.order)
    total_body = sum(len(sec.slides_json or []) for sec in sorted_sections)
    # 표지(1) + 목차(1) + 본문(total_body) + 마무리(1)
    total = 2 + total_body + 1
    footer = brief.footer_text or f"{company.name}  |  {bid.title[:40]}"

    _draw_cover(prs, theme, bid, company)
    _draw_toc(prs, theme, sorted_sections)

    idx = 3
    for sec in sorted_sections:
        idx = _render_section_slides(prs, theme, sec, footer, idx, total)

    closing = _blank(prs)
    _draw_closing(closing, theme, company)

    fname = f"{bid.bid_id}_pt_v{draft.version}.pptx"
    out = OUTPUT_DIR / fname
    prs.save(out)
    logger.info(
        f"📊 PPTX v2 생성: {out} (theme={brief.theme_key}, master={brief.master_label}, 본문 {total_body}장)"
    )
    return out
