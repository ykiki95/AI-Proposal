"""
PPTX 발표자료 생성기 - 다중 슬라이드 타입 + 디자인 템플릿 지원.

지원 타입(슬라이드 dict의 'type'):
- bullets       : 좌측 액센트 + 불릿 본문
- table         : 마크다운/리스트 표
- metric        : 1~3개의 큰 정량 수치 강조
- case_table    : 수행실적 표 (year/client/title/value)
- solution      : 솔루션 카드 (name/tagline/points)
- architecture  : 박스 다이어그램 (모노스페이스)
- section_break : 챕터 구분 페이지 (큰 타이틀 + 부제)
- closing       : 감사·연락처 슬라이드

호환: 과거 [(title, body), ...] 튜플 입력도 그대로 지원 (모두 bullets로 처리).

디자인 템플릿: navy (기본) / mono / warm
- 회사 정보 탭에서 선택하면 db.set_design_template으로 저장되고,
  이 모듈이 db.get_design_template()으로 자동 적용한다.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import List, Optional, Sequence, Union

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import OUTPUT_DIR, get_effective_company

# ---------------------------------------------------------------------------
# 디자인 팔레트
# ---------------------------------------------------------------------------

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x4B, 0x55, 0x63)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


PALETTES = {
    # 공공·SI 정통 네이비 (기본)
    "navy": dict(
        primary=RGBColor(0x0B, 0x2C, 0x5A),
        accent=RGBColor(0x1E, 0x6F, 0xE6),
        soft=RGBColor(0xE6, 0xF0, 0xFB),
        warm=RGBColor(0xF5, 0x8B, 0x00),
    ),
    # 모노 — 컨설팅 스타일
    "mono": dict(
        primary=RGBColor(0x1F, 0x29, 0x37),
        accent=RGBColor(0x4B, 0x55, 0x63),
        soft=RGBColor(0xF3, 0xF4, 0xF6),
        warm=RGBColor(0xC2, 0x41, 0x0C),
    ),
    # 따뜻한 톤 — 민간/스타트업
    "warm": dict(
        primary=RGBColor(0x7C, 0x2D, 0x12),
        accent=RGBColor(0xEA, 0x58, 0x0C),
        soft=RGBColor(0xFE, 0xF3, 0xC7),
        warm=RGBColor(0x16, 0x65, 0x34),
    ),
}


def _palette(template: Optional[str]) -> dict:
    return PALETTES.get((template or "navy").lower(), PALETTES["navy"])


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# 공용 헬퍼
# ---------------------------------------------------------------------------

def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, left, top, width, height, text, *, size=14, bold=False,
          color=GRAY_DARK, align=PP_ALIGN.LEFT, font="맑은 고딕") -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_bar(slide, idx: int, title: str, pal: dict, accent_color: Optional[RGBColor] = None) -> None:
    """본문 슬라이드 공통 상단 — 좌측 액센트 + 옅은 톤 타이틀 영역."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, WHITE)
    bar_color = accent_color or pal["accent"]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SLIDE_H)
    _solid(bar, bar_color)
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.22), 0, SLIDE_W - Inches(0.22), Inches(0.75)
    )
    _solid(title_bg, pal["soft"])
    _text(slide, Inches(0.45), Inches(0.18), Inches(0.7), Inches(0.45),
          f"{idx:02d}", size=16, bold=True, color=bar_color)
    _text(slide, Inches(1.1), Inches(0.2), Inches(11.5), Inches(0.5),
          title, size=18, bold=True, color=pal["primary"])


def _footer(slide, footer: str, idx: int, total: int, pal: dict) -> None:
    foot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3)
    )
    _solid(foot, pal["primary"])
    _text(slide, Inches(0.4), SLIDE_H - Inches(0.28), Inches(10.0), Inches(0.26),
          footer, size=9, color=WHITE)
    _text(slide, SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.28), Inches(1.0), Inches(0.26),
          f"{idx} / {total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# 슬라이드 타입별 렌더러
# ---------------------------------------------------------------------------

def _draw_cover(prs: Presentation, title: str, company, pal: dict) -> None:
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, pal["primary"])
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.7), 0, Inches(1.63), SLIDE_H)
    _solid(accent, pal["accent"])
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.3), Inches(0.12), Inches(2.3)
    )
    _solid(line, pal["warm"])
    _text(slide, Inches(0.85), Inches(2.0), Inches(10.0), Inches(0.5),
          "제 안 서", size=16, bold=True, color=pal["soft"])
    _text(slide, Inches(0.85), Inches(2.55), Inches(10.5), Inches(2.0),
          title, size=32, bold=True, color=WHITE)
    _text(slide, Inches(0.85), Inches(5.5), Inches(10.0), Inches(0.4),
          f"제안사  |  {company.name}", size=14, color=WHITE)
    _text(slide, Inches(0.85), Inches(5.95), Inches(10.0), Inches(0.4),
          f"대표이사  |  {company.ceo}", size=12, color=pal["soft"])
    _text(slide, Inches(0.85), Inches(6.7), Inches(10.0), Inches(0.4),
          f"제출일자  |  {datetime.now().strftime('%Y년 %m월 %d일')}",
          size=11, color=pal["soft"])


def _draw_toc(prs: Presentation, slides: List[dict], pal: dict) -> None:
    slide = _blank(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, WHITE)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.85))
    _solid(band, pal["primary"])
    _text(slide, Inches(0.5), Inches(0.22), Inches(12.0), Inches(0.55),
          "C O N T E N T S", size=22, bold=True, color=WHITE)
    titles = [s.get("title", "") for s in slides]
    half = (len(titles) + 1) // 2
    for col, sub in enumerate([titles[:half], titles[half:]]):
        for i, t in enumerate(sub):
            idx = col * half + i + 1
            x = Inches(0.7 + col * 6.3)
            y = Inches(1.3 + i * 0.42)
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.32), Inches(0.32))
            _solid(circ, pal["accent"])
            tf = circ.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = f"{idx:02d}"
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = WHITE
            _text(slide, x + Inches(0.45), y + Inches(0.02), Inches(5.7), Inches(0.35),
                  t, size=12, bold=True, color=GRAY_DARK)


def _draw_section_break(slide, idx_label: str, title: str, subtitle: str, pal: dict) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, pal["primary"])
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(3.4), Inches(1.2), Inches(0.12)
    )
    _solid(accent, pal["warm"])
    _text(slide, Inches(1.4), Inches(2.7), Inches(11.0), Inches(0.5),
          f"PART  {idx_label}", size=14, bold=True, color=pal["warm"])
    _text(slide, Inches(1.4), Inches(3.3), Inches(11.0), Inches(1.4),
          title, size=40, bold=True, color=WHITE)
    if subtitle:
        _text(slide, Inches(1.4), Inches(4.7), Inches(11.0), Inches(0.7),
              subtitle, size=16, color=pal["soft"])


def _draw_bullets(slide, body: str, pal: dict) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.4), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    lines = [ln.strip().lstrip("-•").strip() for ln in (body or "").split("\n") if ln.strip()]
    if not lines:
        lines = ["(내용 없음)"]
    for i, line in enumerate(lines[:14]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        marker = p.add_run()
        marker.text = "▪  "
        marker.font.size = Pt(13)
        marker.font.color.rgb = pal["warm"]
        marker.font.bold = True
        run = p.add_run()
        run.text = line
        run.font.name = "맑은 고딕"
        run.font.size = Pt(13)
        run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(4)


def _draw_table(slide, table_data, pal: dict) -> None:
    if not table_data:
        return
    if isinstance(table_data, str):
        rows = []
        for raw in table_data.splitlines():
            raw = raw.strip()
            if not raw or set(raw) <= set("-| "):
                continue
            cells = [c.strip() for c in raw.strip("|").split("|")]
            if cells:
                rows.append(cells)
        table_data = rows
    if not table_data:
        return
    rows = len(table_data)
    cols = max(len(r) for r in table_data)
    left, top = Inches(0.55), Inches(1.05)
    width, height = Inches(12.4), Inches(5.9)
    shp = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for ri, row in enumerate(table_data):
        for ci in range(cols):
            cell = shp.cell(ri, ci)
            cell.text = str(row[ci]) if ci < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "맑은 고딕"
                    r.font.size = Pt(11 if ri > 0 else 12)
                    r.font.bold = (ri == 0)
                    r.font.color.rgb = WHITE if ri == 0 else GRAY_DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = pal["primary"] if ri == 0 else (pal["soft"] if ri % 2 == 0 else WHITE)


def _draw_metric(slide, metrics: list, pal: dict) -> None:
    if not metrics:
        metrics = [{"value": "—", "label": "수치 미입력"}]
    metrics = metrics[:4]
    n = len(metrics)
    gap = Inches(0.25)
    card_w = Inches((12.4 - 0.25 * (n - 1)) / n)
    top = Inches(1.8)
    height = Inches(3.6)
    for i, m in enumerate(metrics):
        x = Inches(0.55) + (card_w + gap) * i
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, height)
        _solid(card, pal["soft"])
        _text(slide, x, top + Inches(0.6), card_w, Inches(1.4),
              str(m.get("value", "")), size=44, bold=True, color=pal["primary"], align=PP_ALIGN.CENTER)
        _text(slide, x, top + Inches(2.1), card_w, Inches(0.5),
              str(m.get("label", "")), size=13, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        if m.get("note"):
            _text(slide, x, top + Inches(2.7), card_w, Inches(0.5),
                  str(m["note"]), size=10, color=GRAY, align=PP_ALIGN.CENTER)
    # 하단 보조 설명 영역
    _text(slide, Inches(0.55), Inches(5.7), Inches(12.4), Inches(0.5),
          "※ 위 수치는 자사 보유 자산 또는 RFP 분석 기반 추정치이며, 본문에 근거 출처를 함께 기재합니다.",
          size=10, color=GRAY, align=PP_ALIGN.CENTER)


def _draw_case_table(slide, cases: list, pal: dict) -> None:
    if not cases:
        cases = [{"year": "-", "client": "-", "title": "(등록된 실적 없음)", "value": "-"}]
    rows = [["연도", "발주처", "사업명", "규모/역할"]]
    for c in cases[:10]:
        rows.append([
            str(c.get("year", "")),
            str(c.get("client", "")),
            str(c.get("title", "")),
            str(c.get("value", "")),
        ])
    _draw_table(slide, rows, pal)


def _draw_solution(slide, name: str, tagline: str, points: list, pal: dict) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.55), Inches(1.1), Inches(5.6), Inches(5.7))
    _solid(card, pal["primary"])
    _text(slide, Inches(0.85), Inches(1.4), Inches(5.0), Inches(0.4),
          "OUR SOLUTION", size=11, bold=True, color=pal["warm"])
    _text(slide, Inches(0.85), Inches(1.85), Inches(5.0), Inches(1.0),
          name or "(솔루션명)", size=26, bold=True, color=WHITE)
    if tagline:
        _text(slide, Inches(0.85), Inches(3.0), Inches(5.0), Inches(3.0),
              tagline, size=13, color=pal["soft"])
    _text(slide, Inches(6.5), Inches(1.15), Inches(6.5), Inches(0.4),
          "핵심 포인트", size=14, bold=True, color=pal["primary"])
    box = slide.shapes.add_textbox(Inches(6.5), Inches(1.65), Inches(6.5), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    pts = points or ["(특장점 미등록)"]
    for i, pt in enumerate(pts[:8]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        marker = p.add_run()
        marker.text = "✔  "
        marker.font.size = Pt(13)
        marker.font.color.rgb = pal["warm"]
        marker.font.bold = True
        r = p.add_run()
        r.text = str(pt)
        r.font.name = "맑은 고딕"
        r.font.size = Pt(13)
        r.font.color.rgb = GRAY_DARK
        p.space_after = Pt(4)


def _draw_architecture(slide, body: str, pal: dict) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.55), Inches(1.05), Inches(12.4), Inches(5.9))
    _solid(panel, pal["soft"])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(12.1), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    lines = (body or "(다이어그램 없음)").split("\n")
    for i, line in enumerate(lines[:30]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(11)
        r.font.color.rgb = GRAY_DARK


def _draw_closing(slide, title: str, subtitle: str, company, pal: dict) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, pal["primary"])
    _text(slide, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.4),
          title or "감사합니다", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _text(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
              subtitle, size=18, color=pal["soft"], align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
          f"{company.name}  |  대표 {company.ceo}",
          size=12, color=pal["soft"], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# 디스패처
# ---------------------------------------------------------------------------

SlideInput = Union[dict, tuple]


def _normalize(slides: Sequence[SlideInput]) -> List[dict]:
    out = []
    for s in slides:
        if isinstance(s, dict):
            out.append(s)
        elif isinstance(s, (tuple, list)) and len(s) >= 2:
            out.append({"type": "bullets", "title": s[0], "body": s[1]})
    return out


def _render_one(slide, item: dict, idx: int, total: int, footer: str, company, pal: dict) -> None:
    typ = item.get("type", "bullets")
    title = item.get("title", "")

    if typ == "section_break":
        _draw_section_break(slide, f"{idx:02d}", title, item.get("subtitle", ""), pal)
        return

    if typ == "closing":
        _draw_closing(slide, title, item.get("subtitle", ""), company, pal)
        return

    accent = pal["warm"] if typ == "metric" else pal["accent"]
    _title_bar(slide, idx, title, pal, accent_color=accent)

    if typ == "table":
        _draw_table(slide, item.get("table") or item.get("body", ""), pal)
    elif typ == "metric":
        _draw_metric(slide, item.get("metrics") or [], pal)
    elif typ == "case_table":
        _draw_case_table(slide, item.get("cases") or [], pal)
    elif typ == "solution":
        _draw_solution(slide, item.get("name", ""), item.get("tagline", ""),
                       item.get("points") or [], pal)
    elif typ == "architecture":
        _draw_architecture(slide, item.get("body", ""), pal)
    else:
        _draw_bullets(slide, item.get("body", ""), pal)

    _footer(slide, footer, idx, total, pal)


def render_pptx(bid_id: str, title: str, slides: Sequence[SlideInput],
                template: Optional[str] = None) -> Path:
    """최종 PPTX를 생성하고 경로 반환.
    template: 'navy' | 'mono' | 'warm'. None이면 DB 설정값 적용."""
    if template is None:
        try:
            from tools.db import get_design_template
            template = get_design_template()
        except Exception:
            template = "navy"
    pal = _palette(template)
    company = get_effective_company()
    items = _normalize(slides)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    footer = f"{company.name}  |  {title[:40]}"
    total = len(items)

    _draw_cover(prs, title, company, pal)
    toc_items = [s for s in items if s.get("type") not in ("section_break", "closing")]
    if toc_items:
        _draw_toc(prs, toc_items, pal)

    for i, item in enumerate(items, start=1):
        slide = _blank(prs)
        _render_one(slide, item, i, total, footer, company, pal)

    fname = f"{bid_id}_pt.pptx"
    out_path = OUTPUT_DIR / fname
    prs.save(out_path)
    logger.info(f"PPTX 생성: {out_path} (template={template}, 본문 {total}장)")
    return out_path
