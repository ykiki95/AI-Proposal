"""기본 PPTX 마스터 3종 자동 생성기.

회사 정보 탭에서 '🎨 기본 마스터 3종 자동 생성' 버튼을 누르면 호출된다.
3종 테마(corporate_navy / innovation_orange / minimal_white)에 대해
표지·간지·본문·2단·메트릭 5장이 든 미리보기 .pptx를 만들어 storage/masters/ 에 저장하고,
회사 정보 탭의 PPTX 마스터 4슬롯 중 빈 슬롯에 자동 등록한다.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Dict

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import STORAGE_DIR
from tools import db
from tools.pptx_themes import THEMES, Theme

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, x, y, w, h, text, *, size=14, bold=False, color=None,
          align=PP_ALIGN.LEFT, font="맑은 고딕") -> None:
    if color is None:
        color = RGBColor(0x1F, 0x29, 0x37)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _draw_cover(prs, theme: Theme) -> None:
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, theme.rgb("primary"))
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.7), 0, Inches(1.63), SLIDE_H)
    _solid(accent, theme.rgb("accent"))
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.3), Inches(0.12), Inches(2.3))
    _solid(line, theme.rgb("warm"))
    _text(s, Inches(0.85), Inches(2.0), Inches(10.0), Inches(0.5),
          "MASTER  PREVIEW", size=14, bold=True, color=theme.rgb("soft"), font=theme.font_title)
    _text(s, Inches(0.85), Inches(2.55), Inches(11.0), Inches(2.0),
          theme.label, size=36, bold=True, color=WHITE, font=theme.font_title)
    _text(s, Inches(0.85), Inches(5.5), Inches(10.0), Inches(0.5),
          theme.tone, size=14, color=theme.rgb("soft"), font=theme.font_body)
    _text(s, Inches(0.85), Inches(6.7), Inches(10.0), Inches(0.4),
          f"생성일 {datetime.now().strftime('%Y-%m-%d')}",
          size=11, color=theme.rgb("soft"), font=theme.font_body)


def _draw_section_divider(prs, theme: Theme, label: str) -> None:
    s = _blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, theme.rgb("primary"))
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.4), Inches(1.2), Inches(0.12))
    _solid(line, theme.rgb("warm"))
    _text(s, Inches(1.4), Inches(2.7), Inches(11.0), Inches(0.5),
          "PART  01", size=14, bold=True, color=theme.rgb("warm"), font=theme.font_title)
    _text(s, Inches(1.4), Inches(3.3), Inches(11.0), Inches(1.4),
          label, size=44, bold=True, color=WHITE, font=theme.font_title)
    _text(s, Inches(1.4), Inches(4.8), Inches(11.0), Inches(0.6),
          "간지(章) 슬라이드 — 큰 챕터 구분", size=16, color=theme.rgb("soft"), font=theme.font_body)


def _title_bar(s, theme: Theme, idx: int, title: str) -> None:
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _solid(bg, WHITE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), SLIDE_H)
    _solid(bar, theme.rgb("accent"))
    title_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.22), 0, SLIDE_W - Inches(0.22), Inches(0.75))
    _solid(title_bg, theme.rgb("soft"))
    _text(s, Inches(0.45), Inches(0.18), Inches(0.7), Inches(0.45),
          f"{idx:02d}", size=16, bold=True, color=theme.rgb("accent"), font=theme.font_title)
    _text(s, Inches(1.1), Inches(0.2), Inches(11.5), Inches(0.5),
          title, size=18, bold=True, color=theme.rgb("primary"), font=theme.font_title)


def _draw_body(prs, theme: Theme) -> None:
    s = _blank(prs)
    _title_bar(s, theme, 1, "본문(불릿) 페이지 — 제목 + 5~7개 불릿")
    box = s.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(12.4), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    sample = [
        "본문 첫 줄 — 핵심 결론을 두괄식으로",
        "Why → How → Effect 3박자 서사",
        "발주처 RFP 단어를 의도적으로 그대로 인용",
        "정량 수치(예: 90.05%, -45%) 를 굵게 강조",
        "차별화 포인트 1줄 별도 단락",
        "다음 슬라이드 예고로 마무리",
    ]
    for i, line in enumerate(sample):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        m = p.add_run()
        m.text = "▪  "
        m.font.size = Pt(13)
        m.font.color.rgb = theme.rgb("warm")
        m.font.bold = True
        r = p.add_run()
        r.text = line
        r.font.name = theme.font_body
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
        p.space_after = Pt(6)


def _draw_two_column(prs, theme: Theme) -> None:
    s = _blank(prs)
    _title_bar(s, theme, 2, "2단 비교 — AS-IS vs TO-BE")
    col_w = Inches(6.0)
    col_h = Inches(5.6)
    top = Inches(1.1)
    for i, (label, color, items) in enumerate([
        ("AS-IS  현행", theme.rgb("text_gray"),
         ["수기 처리 / 종이 보고", "데이터 단절", "현장 활용 어려움"]),
        ("TO-BE  개선 후", theme.rgb("accent"),
         ["AI 자동 분류·요약", "통합 데이터 허브", "현장 모바일 즉시 활용"]),
    ]):
        x = Inches(0.55) + (col_w + Inches(0.3)) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, col_w, col_h)
        _solid(card, theme.rgb("soft") if i == 0 else WHITE)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, col_w, Inches(0.55))
        _solid(head, color)
        _text(s, x + Inches(0.25), top + Inches(0.1), col_w, Inches(0.4),
              label, size=14, bold=True, color=WHITE, font=theme.font_title)
        box = s.shapes.add_textbox(x + Inches(0.3), top + Inches(0.8),
                                   col_w - Inches(0.6), col_h - Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            m = p.add_run()
            m.text = "•  "
            m.font.color.rgb = color
            m.font.bold = True
            r = p.add_run()
            r.text = it
            r.font.name = theme.font_body
            r.font.size = Pt(13)
            r.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
            p.space_after = Pt(6)


def _draw_metric(prs, theme: Theme) -> None:
    s = _blank(prs)
    _title_bar(s, theme, 3, "정량 카드 — 핵심 수치 강조")
    metrics = [("90.05%", "음성인식 정확도"), ("-45%", "행정 처리시간"), ("3,200건", "월간 처리량")]
    n = len(metrics)
    gap = Inches(0.25)
    card_w = Inches((12.4 - 0.25 * (n - 1)) / n)
    top = Inches(1.8)
    h = Inches(3.6)
    for i, (val, lbl) in enumerate(metrics):
        x = Inches(0.55) + (card_w + gap) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, h)
        _solid(card, theme.rgb("soft"))
        _text(s, x, top + Inches(0.7), card_w, Inches(1.5),
              val, size=48, bold=True, color=theme.rgb("primary"),
              align=PP_ALIGN.CENTER, font=theme.font_title)
        _text(s, x, top + Inches(2.3), card_w, Inches(0.5),
              lbl, size=14, color=RGBColor(0x4B, 0x55, 0x63),
              align=PP_ALIGN.CENTER, font=theme.font_body)


def build_master_preview(theme: Theme, out_dir: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    _draw_cover(prs, theme)
    _draw_section_divider(prs, theme, "사업 이해")
    _draw_body(prs, theme)
    _draw_two_column(prs, theme)
    _draw_metric(prs, theme)
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f"master_{theme.key}.pptx"
    prs.save(out)
    logger.info(f"마스터 미리보기 생성: {out}")
    return out


def build_all_master_previews_and_register(masters_dir: Path | None = None) -> Dict[str, Path]:
    """3종 마스터 미리보기 PPTX 생성 + 회사 정보 4슬롯에 자동 등록.

    이미 등록된 슬롯은 건드리지 않는다. 빈 슬롯이 부족하면 가능한 만큼만 등록.
    """
    masters_dir = masters_dir or (STORAGE_DIR / "masters")
    profile = db.get_company_profile() or {}
    out: Dict[str, Path] = {}
    used_slots = set()
    for i in range(1, 5):
        if profile.get(f"pptx_master_{i}_path"):
            used_slots.add(i)

    next_slot = next((i for i in range(1, 5) if i not in used_slots), None)
    save_payload: Dict = {}
    for theme_key, theme in THEMES.items():
        path = build_master_preview(theme, masters_dir)
        out[theme_key] = path
        if next_slot is not None:
            save_payload[f"pptx_master_{next_slot}_path"] = str(path)
            save_payload[f"pptx_master_{next_slot}_label"] = f"기본 · {theme.label}"
            used_slots.add(next_slot)
            next_slot = next((i for i in range(1, 5) if i not in used_slots), None)
    if save_payload:
        db.save_company_profile(save_payload)
    return out
